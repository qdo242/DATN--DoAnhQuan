#!/usr/bin/env python3
"""Create PowerPoint slides for thesis defense."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), '..', '..', 'doandocs')

SLIDES = [
    {"type": "title", "title": "XAY DUNG GIAI PHAP TRUYN TIN BAO MAT GIUA CAC THIET BI IOT", "subtitle": "Do Anh Quan\nGV: Thay Hai"},
    {"type": "content", "title": "MUC LUC", "body": "Chuong 1: Gioi thieu\nChuong 2: Ngon ngu\nChuong 3: Thiet ke\nChuong 4: Trien khai\nChuong 5: Kiem thu"},
    {"type": "content", "title": "CHUONG 1: GIOI THIEU", "body": "1.1 Boi canh\n1.2 Muc tieu\n1.3 Pham vi"},
    {"type": "content", "title": "CHUONG 2: NGON NGU", "body": "2.1 ESP32\n2.2 AES-128-CBC\n2.3 Flask\n2.4 SQLite"},
    {"type": "content", "title": "CHUONG 3: THIET KE", "body": "3.1 Kien truc\n3.2 Luong du lieu\n3.3 Protocol\n3.4 Database"},
    {"type": "content", "title": "CHUONG 4: TRIEN KHAI", "body": "4.1 Wokwi ESP32\n4.2 Server Flask\n4.3 AES Encrypt\n4.4 Anti-replay\n4.5 Android App\n4.6 Benchmark\n4.7 Test"},
    {"type": "content", "title": "CHUONG 5: KIEM THU", "body": "5.1 Muc tieu\n5.2 Test case\n5.2 Ket qua"},
    {"type": "content", "title": "KET LUAN", "body": "Da hoan thanh:\n- AES-128-CBC encryption\n- Anti-replay protection\n- Benchmark AES vs XOR"},
    {"type": "content", "title": "CAM ON", "body": "Loi on thay\nLoi on ban be"},
]

def create_slides():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for slide_data in SLIDES:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = slide_data["title"]
        
        if slide_data["type"] == "title":
            subtitle = slide.placeholders[1]
            subtitle.text = slide_data.get("subtitle", "")
        else:
            body = slide.placeholders[1]
            body.text = slide_data.get("body", "")
    
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    output_path = os.path.join(OUTPUT_DIR, "Slide_DATN_HUCE.pptx")
    prs.save(output_path)
    print(f"Slides saved to: {output_path}")

if __name__ == '__main__':
    create_slides()